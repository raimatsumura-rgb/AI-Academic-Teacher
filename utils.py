import io
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import streamlit as st

def extract_text_from_files(uploaded_files):
    """
    دالة خارقة لاستخراج النصوص من PDF, DOCX, PPTX, TXT
    تدعم ترميز UTF-8 لضمان سلامة اللغة العربية واليابانية.
    """
    all_text = ""
    
    for uploaded_file in uploaded_files:
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        try:
            # 1. معالجة ملفات PDF
            if file_extension == 'pdf':
                pdf_reader = PdfReader(uploaded_file)
                for page in pdf_reader.pages:
                    content = page.extract_text()
                    if content:
                        all_text += content + "\n"
            
            # 2. معالجة ملفات Word (DOCX)
            elif file_extension == 'docx':
                doc = DocxDocument(uploaded_file)
                for para in doc.paragraphs:
                    all_text += para.text + "\n"
            
            # 3. معالجة ملفات PowerPoint (PPTX)
            elif file_extension == 'pptx':
                ppt = Presentation(uploaded_file)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            all_text += shape.text + "\n"
            
            # 4. معالجة ملفات النصوص (TXT)
            elif file_extension == 'txt':
                # قراءة الملف مع دعم ترميز UTF-8 للغات العالمية
                stringio = io.StringIO(uploaded_file.getvalue().decode("utf-8"))
                all_text += stringio.read() + "\n"
                
        except Exception as e:
            st.error(f"Error processing {uploaded_file.name}: {e}")
            
    return all_text